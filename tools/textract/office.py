#!/usr/bin/env python
# -*- coding:utf-8 -*-
import os
import pandas as pd
import pdfplumber
from docx import Document
from pptx import Presentation
import win32com.client
import pythoncom

"""
pip install python-docx==0.8.11
pip install python-pptx==0.6.21
pip install pdfplumber==0.6.2
"""

# import textract
# text = textract.process(r'E:\tmp\zdump\test.pptx')
# print(text.decode('utf-8'))       # 无法获取表格内容


def docx(filepath):
    texts = list()
    fopen = Document(filepath)
    for p in fopen.paragraphs:
        texts.append(p.text)
    for table in fopen.tables:
        for row in table.rows:
            for cell in row.cells:
                texts.append(cell.text)
    return set(texts)


def pptx(filepath):
    texts = list()
    fopen = Presentation(filepath)
    for slide in fopen.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)
            if shape.has_table:
                tbl = shape.table
                for r in range(0, len(tbl.rows)):
                    for c in range(0, len(tbl.columns)):
                        cell = tbl.cell(r, c)
                        texts.append(cell.text_frame.text)
    return set(texts)


def xlsx(filepath):
    texts = list()
    fopen = pd.read_excel(filepath, sheet_name=None)
    for name, sheet in fopen.items():
        for index, row in sheet.iterrows():
            for value in row:
                texts.append(str(value))
    return set(texts)


def doc(filepath):
    # https://www.cnblogs.com/zhuminghui/p/11765401.html
    texts = list()
    pythoncom.CoInitialize()
    app = win32com.client.Dispatch('Word.Application')
    # 如果不声明下列属性，运行的时候会显示的打开office软件操作文档
    app.Visible = 0            # 后台运行
    app.DisplayAlerts = 0      # 不显示，不警告
    fopen = app.Documents.Open(filepath)
    # python3.5.4 + pywin32-301提取失败
    # for paragraph in fopen.Paragraphs:
    #     texts.append(paragraph.Range.Text)
    # for table in fopen.Tables:
    #     for row in table.Rows:
    #         for cell in row.Cells:
    #             texts.append(cell.Range.Text)
    new_pdf_file = filepath + '.pdf'
    # https://stackoverflow.com/questions/30481136/pywin32-save-docx-as-pdf
    fopen.SaveAs(new_pdf_file, FileFormat=17)
    fopen.Close()
    app.Quit()
    texts = pdf(new_pdf_file)
    os.remove(new_pdf_file)
    pythoncom.CoUninitialize()
    return set(texts)


def ppt(filepath):
    """
    线程中使用win32com发生错误: pywintypes.com_error: (-2147221008, '尚未调用 CoInitialize。', None, None)
    解决方案：    # 线程初始化
                pythoncom.CoInitialize()
                # 释放资源
                pythoncom.CoUninitialize()
    """
    pythoncom.CoInitialize()
    app = win32com.client.DispatchEx('Powerpoint.Application')
    ppt.DisplayAlerts = 0                   # 不显示,不警告
    fopen = app.Presentations.Open(filepath)
    new_pdf_file = filepath + '.pdf'
    # fopen.SaveAs(new_pptx_file, 24)               # 保存为pptx
    # https://community.esri.com/t5/python-questions/how-can-i-use-python-to-convert-ppt-to-pdf/td-p/310114
    fopen.SaveAs(new_pdf_file, FileFormat=32)       # 保存为pdf
    fopen.Close()
    app.Quit()
    texts = pdf(new_pdf_file)
    os.remove(new_pdf_file)
    pythoncom.CoUninitialize()
    return texts


# def ppt(filepath):
#     ppt = win32com.client.DispatchEx('Powerpoint.Application')
#     ppt.DisplayAlerts = 0  # 不显示,不警告
#     file = ppt.Presentations.Open(filepath)
#     # 用pywin32实在不知道怎么访问ppt文本和表格内容
#     for slide in file.Slides:
#         for shape in slide.Shapes:
#             # https://docs.microsoft.com/zh-cn/office/vba/api/powerpoint.shape.textframe
#             # https://docs.microsoft.com/zh-cn/office/vba/api/powerpoint.textrange
#             # https://docs.microsoft.com/zh-cn/office/vba/api/powerpoint.textrange.text
#             print(shape.TextFrame.TextRange)                        # 任意文本框里的内容
#             # Exception: pywintypes.com_error: (-2147352567, '发生意外。', (0, None, None, None, 0, -2147467259), None)
#             # import win32api
#             # win32api.FormatMessage(-2147467259)
#             # >> '未指定的错误\r\n'
#             # https://blog.csdn.net/huicheng0703/article/details/122242972
#             # https://www.cnblogs.com/vhills/p/8098715.html
#             # for paragraph in shape.TextFrame.TextRange.Paragraphs:
#             #     print(paragraph.Text)


def xls(filepath):
    return xlsx(filepath)


def pdf(filepath):
    texts = list()
    with pdfplumber.open(filepath) as fopen:
        for page in fopen.pages:
            text = page.extract_text()  # 提取文本
            if text:
                texts.append(str(text))
    return set(texts)


__office_methods = {
    'doc': doc,
    'docx': docx,
    'xls': xls,
    'xlsx': xlsx,
    'ppt': ppt,
    'pptx': pptx,
    'pdf': pdf,
}


def autocheck(filepath):
    suffix = os.path.basename(filepath).split('.')[-1].lower()
    if suffix in __office_methods:
        return '\n'.join(__office_methods[suffix](filepath))
    return ''

